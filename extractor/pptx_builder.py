"""
PPTX export builder — programmatic slide generation from deal records.

Reads brand.json for color consistency, builds 16:9 slides with editable text shapes.
"""
import json
import os
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None

try:
    import requests
except ImportError:
    requests = None


def build_deck(
    deals: List[Dict[str, Any]],
    include_pipeline_summary: bool = False,
    mapbox_token: Optional[str] = None,
) -> bytes:
    """
    Build a PPTX deck from deal records.

    Args:
        deals: List of deal record dicts (extracted_fields, showcase, underwrite.runs)
        include_pipeline_summary: Add a summary table slide
        mapbox_token: Mapbox API token for static map images

    Returns:
        PPTX file bytes
    """
    if Presentation is None:
        raise ImportError("python-pptx not installed")

    # Load brand colors
    brand = _load_brand()

    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Add optional pipeline summary slide
    if include_pipeline_summary and deals:
        _add_pipeline_summary_slide(prs, deals, brand)

    # Add deal slides
    for deal in deals:
        _add_deal_overview_slide(prs, deal, brand, mapbox_token)

        # Add financial overview if passed run with CFO exists
        if _get_latest_passed_run(deal):
            _add_financial_overview_slide(prs, deal, brand)

    # Convert to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def _load_brand() -> Dict[str, str]:
    """Load brand.json color tokens."""
    try:
        path = Path(__file__).parent / "brand.json"
        if path.exists():
            with open(path) as f:
                return json.load(f)
    except Exception:
        pass
    # Fallback
    return {
        "purple": "#7D5A7D",
        "purpleLight": "#E6DCE6",
        "ink": "#1F1F1F",
        "cardBg": "#F5F4F6",
    }


def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def _add_pipeline_summary_slide(
    prs: "Presentation", deals: List[Dict[str, Any]], brand: Dict[str, str]
) -> None:
    """Add a summary table slide."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Pipeline Summary"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))

    # Add table (simple text-based approximation)
    table_shape = slide.shapes.add_table(
        rows=len(deals) + 1, cols=10, left=Inches(0.5), top=Inches(1.1), width=Inches(9), height=Inches(4)
    )
    table = table_shape.table

    # Header row
    headers = ["Asset", "Market", "Tenants", "Occupancy", "NIY", "RY", "On/Off", "Price £m", "WAULT", "Comment"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.bold = True

    # Data rows (zebra)
    for row_idx, deal in enumerate(deals, start=1):
        fields = deal.get("extracted_fields", {})
        asset = fields.get("Project Name", "—")
        market = (deal.get("market_ids") or ["—"])[0]
        tenants = fields.get("Number of Tenants", "—")
        occupancy = fields.get("Economic occupancy rate, %", "—")
        niy = fields.get("Yield", "—")
        ry = fields.get("Yield2", "—")
        price = fields.get("Deal value, CCY", "—")
        wault = fields.get("WAULT, years", "—")
        comment = fields.get("Comment", "")[:20]

        values = [asset, market, str(tenants), f"{occupancy}%", f"{niy}%", f"{ry}%", "Active", f"£{price / 1_000_000:.1f}" if isinstance(price, (int, float)) else "—", str(wault), comment]

        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 244, 246)  # cardBg light
            cell.text_frame.paragraphs[0].font.size = Pt(9)


def _add_deal_overview_slide(
    prs: "Presentation", deal: Dict[str, Any], brand: Dict[str, str], mapbox_token: Optional[str]
) -> None:
    """Add deal overview slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))  # Rectangle
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))
    title_shape.line.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))

    fields = deal.get("extracted_fields", {})
    asset = fields.get("Project Name", "Asset")
    location = deal.get("showcase", {}).get("location", {}).get("address", "—")
    price = fields.get("Deal value, CCY")
    ry = fields.get("Yield2")

    title_text = f"{asset} | {location} | £{price / 1_000_000:.1f}m | {ry}%" if isinstance(price, (int, float)) else f"{asset} | {location} | TBC"
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # KPI grid (simplified text version)
    showcase = deal.get("showcase") or {}
    kpis = showcase.get("kpis") or {}

    left = Inches(0.5)
    top = Inches(1.0)
    _add_kpi_grid(slide, kpis, left, top, brand)

    # Rationale bullets
    left = Inches(0.5)
    top = Inches(3.0)
    _add_section_band(slide, "Investment Rationale", left, top, Inches(5), brand)
    rationale = showcase.get("rationale_bullets", [])
    for idx, bullet in enumerate(rationale[:3]):
        bullet_top = top + Inches(0.5) + Inches(idx * 0.3)
        bullet_box = slide.shapes.add_textbox(left, bullet_top, Inches(4.5), Inches(0.25))
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        bullet_frame.text = f"{bullet.get('label', '')} — {bullet.get('text', '')}"
        bullet_frame.paragraphs[0].font.size = Pt(9)

    # Map placeholder (right side)
    map_left = Inches(5.3)
    map_top = Inches(1.0)
    _add_map_placeholder(slide, deal, map_left, map_top, Inches(4.2), Inches(2.3), brand, mapbox_token)

    # Photo placeholder (right side)
    photo_left = Inches(5.3)
    photo_top = Inches(3.5)
    _add_photo_placeholder(slide, deal, photo_left, photo_top, Inches(4.2), Inches(1.8), brand)


def _add_financial_overview_slide(prs: "Presentation", deal: Dict[str, Any], brand: Dict[str, str]) -> None:
    """Add financial overview slide with CFO data."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    fields = deal.get("extracted_fields", {})
    asset = fields.get("Project Name", "Asset")

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{asset} — Financial Overview"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True

    latest_run = _get_latest_passed_run(deal)
    if not latest_run:
        return

    cfo = latest_run.get("returns", {}).get("cfo")
    if not cfo:
        return

    assumptions = cfo.get("assumptions_table", {})
    unlev = cfo.get("unlevered", {})
    lev = cfo.get("levered", {})

    # Assumptions table
    top = Inches(1.0)
    _add_assumptions_table(slide, assumptions, top, brand)

    # Returns panels
    top = Inches(2.5)
    _add_returns_panels(slide, unlev, lev, top, brand)


def _add_section_band(slide: "Slide", title: str, left, top, width, brand: Dict[str, str]) -> None:
    """Add a section band header."""
    band_shape = slide.shapes.add_shape(1, left, top, width, Inches(0.35))
    band_shape.fill.solid()
    band_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))
    band_shape.line.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))

    band_text = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.25))
    band_frame = band_text.text_frame
    band_frame.text = title
    band_frame.paragraphs[0].font.size = Pt(11)
    band_frame.paragraphs[0].font.bold = True
    band_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def _add_kpi_grid(slide: "Slide", kpis: Dict[str, Any], left, top, brand: Dict[str, str]) -> None:
    """Add KPI grid (2x4)."""
    kpi_items = [
        ("Tenure", kpis.get("tenure")),
        ("Units", kpis.get("units")),
        ("Area (sqft)", kpis.get("lettable_area_sqft")),
        ("Occupancy %", kpis.get("occupancy_pct")),
        ("Passing Rent (psf)", kpis.get("passing_rent_psf")),
        ("Cap Value (psf)", kpis.get("capital_value_psf")),
        ("NIY %", kpis.get("niy_pct")),
        ("RY %", kpis.get("ry_pct")),
    ]

    for idx, (label, value) in enumerate(kpi_items):
        col = idx % 4
        row = idx // 4
        cell_left = left + Inches(col * 1.25)
        cell_top = top + Inches(row * 0.6)

        # Cell background
        cell_shape = slide.shapes.add_shape(1, cell_left, cell_top, Inches(1.2), Inches(0.55))
        cell_shape.fill.solid()
        cell_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(brand["cardBg"]))
        cell_shape.line.color.rgb = RGBColor(200, 200, 200)

        # Label
        label_box = slide.shapes.add_textbox(cell_left + Inches(0.05), cell_top + Inches(0.05), Inches(1.1), Inches(0.2))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(8)
        label_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

        # Value
        value_box = slide.shapes.add_textbox(cell_left + Inches(0.05), cell_top + Inches(0.25), Inches(1.1), Inches(0.25))
        value_frame = value_box.text_frame
        value_frame.text = str(value or "TBC")
        value_frame.paragraphs[0].font.size = Pt(10)
        value_frame.paragraphs[0].font.bold = True
        value_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))


def _add_map_placeholder(slide: "Slide", deal, left, top, width, height, brand, mapbox_token) -> None:
    """Add map placeholder."""
    map_shape = slide.shapes.add_shape(1, left, top, width, height)
    map_shape.fill.solid()
    map_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    map_shape.line.color.rgb = RGBColor(100, 100, 100)

    map_text = slide.shapes.add_textbox(left, top + height / 2 - Inches(0.1), width, Inches(0.2))
    map_frame = map_text.text_frame
    map_frame.text = "Map TBC"
    map_frame.paragraphs[0].font.size = Pt(12)
    map_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    map_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


def _add_photo_placeholder(slide: "Slide", deal, left, top, width, height, brand) -> None:
    """Add photo placeholder."""
    photo_shape = slide.shapes.add_shape(1, left, top, width, height)
    photo_shape.fill.solid()
    photo_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(brand["cardBg"]))
    photo_shape.line.color.rgb = RGBColor(150, 150, 150)

    photo_text = slide.shapes.add_textbox(left, top + height / 2 - Inches(0.1), width, Inches(0.2))
    photo_frame = photo_text.text_frame
    photo_frame.text = "Asset Photo"
    photo_frame.paragraphs[0].font.size = Pt(11)
    photo_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    photo_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


def _add_assumptions_table(slide: "Slide", assumptions: Dict[str, Any], top, brand) -> None:
    """Add key assumptions table."""
    # Simplified text version
    title_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(4), Inches(0.3))
    title_frame = title_box.text_frame
    title_frame.text = "Key Assumptions"
    title_frame.paragraphs[0].font.size = Pt(12)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))


def _add_returns_panels(slide: "Slide", unlev: Dict[str, Any], lev: Dict[str, Any], top, brand) -> None:
    """Add returns side-by-side panels."""
    # Left panel
    left_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(4.5), Inches(2.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = f"Unlevered Returns\n\nIRR: {unlev.get('irr', 0) * 100:.2f}%\nEM: {unlev.get('em', 0):.2f}x\nProfit: £{unlev.get('profit', 0):,.0f}"
    left_frame.paragraphs[0].font.size = Pt(12)
    left_frame.paragraphs[0].font.bold = True
    left_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))

    # Right panel
    right_box = slide.shapes.add_textbox(Inches(5.2), top, Inches(4.5), Inches(2.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.text = f"Levered Returns\n\nIRR: {lev.get('irr', 0) * 100:.2f}%\nEM: {lev.get('em', 0):.2f}x\nProfit: £{lev.get('profit', 0):,.0f}"
    right_frame.paragraphs[0].font.size = Pt(12)
    right_frame.paragraphs[0].font.bold = True
    right_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(brand["purple"]))


def _get_latest_passed_run(deal: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Find the latest passed run with CFO data."""
    underwrite = deal.get("underwrite")
    if not underwrite:
        return None
    runs = underwrite.get("runs", [])
    for run in reversed(runs):
        if run.get("checks", {}).get("pass") and run.get("returns", {}).get("cfo"):
            return run
    return None
